import streamlit as st
import google.generativeai as genai
import PyPDF2
from pptx import Presentation
import json
import os
import asyncio
import edge_tts
from datetime import datetime
from dotenv import load_dotenv
import re

# --- CONFIGURAZIONE SICUREZZA ---
load_dotenv()
API_KEY = os.getenv("GOOGLE_API_KEY")

st.set_page_config(page_title="NEXUS VANGUARD", page_icon="⚡", layout="centered")

# --- FUNZIONI UTILITY ---
def clean_markdown_for_audio(text):
    # Rimuove grassetti (**), corsivi (*) e titoli (##) per la lettura audio
    text = text.replace("**", "").replace("##", "").replace("* ", "")
    return text

async def generate_audio(text, output_filename):
    # Voce maschile italiana profonda e professionale (Diego)
    voice = "it-IT-DiegoNeural" 
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(output_filename)

def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# --- IL CERVELLO ---
def analyze_content(text):
    if not API_KEY:
        raise ValueError("API Key non trovata nel file .env")
        
    genai.configure(api_key=API_KEY)
    # Usiamo l'alias stabile
    model = genai.GenerativeModel('gemini-flash-latest') 

    prompt = f"""
    Analizza questo materiale didattico agendo come un consulente esperto che deve spiegare i concetti a un manager, o uno studente lavoratore brillante che vuole andare al sodo.
    
    Il tuo obiettivo NON è sembrare un libro stampato, ma dimostrare intelligenza pratica.
    Usa un linguaggio professionale ma FLUIDO, NATURALE e PARLATO.

    Restituisci ESCLUSIVAMENTE un output strutturato così:

    1. **5 CONCETTI "TAKEAWAY" (Il Succo del Discorso):**
       Estrai i 5 punti fondamentali. Spiegali in modo diretto.

    2. **3 DOMANDE "DAL MONDO REALE" (Smart Questions):**
       Genera 3 domande da porre al docente.
       Usa formule come "Professore, ma nel caso pratico in cui...", "Come si concilia questo con...".
       Collega la teoria alla pratica operativa (costi, tempi, realtà aziendale).

    3. **3 PUNTI DI ATTENZIONE (Critical Reality Check):**
       Identifica 3 aspetti dove la teoria presentata potrebbe scontrarsi con la dura realtà aziendale.

    --- MATERIALE DA ANALIZZARE ---
    {text[:40000]} 
    """
    
    response = model.generate_content(prompt)
    return response.text

# --- SALVATAGGIO ---
def save_mission(target_date, content, filename_orig, audio_filename):
    filename = "mission_orders.json"
    
    if os.path.exists(filename):
        with open(filename, "r") as f:
            try:
                missions = json.load(f)
            except:
                missions = []
    else:
        missions = []
    
    new_mission = {
        "id": datetime.now().strftime("%Y%m%d%H%M%S"),
        "target_date": target_date.strftime("%Y-%m-%d"),
        "original_file": filename_orig,
        "payload": content,
        "audio_file": audio_filename, # Salviamo il nome del file audio
        "status": "ARMED"
    }
    
    missions.append(new_mission)
    
    with open(filename, "w") as f:
        json.dump(missions, f, indent=4)

# --- INTERFACCIA ---
st.title("⚡ NEXUS VANGUARD")
st.subheader("Tactical Briefing & Audio Generator")

if not API_KEY:
    st.error("❌ API Key mancante nel file .env")
    st.stop()
else:
    st.success("🛡️ SISTEMA ARMATO E CONNESSO (Secure Mode)")

uploaded_file = st.file_uploader("Carica File Target (PDF/PPTX)", type=["pdf", "pptx"])
target_date = st.date_input("Data dell'Ingaggio (Lezione)")

if st.button("ARMA IL SISTEMA"):
    if not uploaded_file:
        st.error("Nessun materiale caricato.")
    else:
        with st.spinner("Analisi tattica + Sintesi Vocale in corso..."):
            try:
                # 1. Estrazione testo
                text_content = ""
                if uploaded_file.name.endswith(".pdf"):
                    text_content = extract_text_from_pdf(uploaded_file)
                elif uploaded_file.name.endswith(".pptx"):
                    text_content = extract_text_from_pptx(uploaded_file)
                
                # 2. Analisi AI
                briefing = analyze_content(text_content)
                
                # 3. Generazione Audio
                audio_filename = f"briefing_{target_date.strftime('%Y-%m-%d')}.mp3"
                clean_text = clean_markdown_for_audio(briefing)
                asyncio.run(generate_audio(clean_text, audio_filename))
                
                # 4. Salvataggio
                save_mission(target_date, briefing, uploaded_file.name, audio_filename)
                
                st.success(f"✅ Missione programmata per il {target_date.strftime('%d/%m/%Y')}")
                st.info(f"🎙️ Audio generato: {audio_filename}")
                
                with st.expander("👁️ Anteprima Testo"):
                    st.markdown(briefing)
                
            except Exception as e:
                st.error(f"Errore: {e}")